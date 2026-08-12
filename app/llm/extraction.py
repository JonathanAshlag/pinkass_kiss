"""Document text/image extraction for the ingestion pipeline.

Dispatches across supported document formats (PDF, DOCX, XLSX, PPTX, HTML,
plain text) and normalizes each into the ``ExtractedContent`` shape consumed
by ``app.llm.pipeline.run_ingestion_pipeline`` (``text`` plus a
``content_parts``-style list of ``{"type": "text", ...}`` /
``{"type": "image_url", ...}`` blocks).
"""

import base64
from dataclasses import dataclass, field
from io import BytesIO

MAX_IMAGES = 20


@dataclass
class ExtractedContent:
    text: str
    parts: list[dict] = field(default_factory=list)


def _encode_image(data: bytes, name: str) -> str:
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else "png"
    mime = f"image/{ext}" if ext in ("png", "jpeg", "jpg", "gif", "webp") else "image/png"
    b64 = base64.b64encode(data).decode("ascii")
    return f"data:{mime};base64,{b64}"


def _append_image(parts: list[dict], image_count: int, data: bytes, name: str) -> int:
    """Encode an image and append it to ``parts`` if under the MAX_IMAGES cap.

    Returns the updated image_count. Shared by the PDF, DOCX, and PPTX
    branches, which each walk a different library's structures but need the
    same capped image_url-appending behavior.
    """
    if image_count >= MAX_IMAGES:
        return image_count
    uri = _encode_image(data, name)
    parts.append({"type": "image_url", "image_url": {"url": uri}})
    return image_count + 1


def _extract_pdf(content: bytes) -> ExtractedContent:
    from pypdf import PdfReader
    reader = PdfReader(BytesIO(content))
    text_parts: list[str] = []
    parts: list[dict] = []
    image_count = 0
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            text_parts.append(page_text)
            parts.append({"type": "text", "text": page_text})
        if image_count < MAX_IMAGES:
            for img in page.images:
                if image_count >= MAX_IMAGES:
                    break
                image_count = _append_image(parts, image_count, img.data, img.name)
    return ExtractedContent(text="\n".join(text_parts), parts=parts)


def _extract_docx(content: bytes) -> ExtractedContent:
    from docx import Document
    doc = Document(BytesIO(content))
    text_parts: list[str] = []
    parts: list[dict] = []
    image_count = 0
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
            parts.append({"type": "text", "text": para.text})
        if image_count < MAX_IMAGES:
            for run in para.runs:
                for drawing in run.element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                    blips = drawing.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    for blip in blips:
                        if image_count >= MAX_IMAGES:
                            break
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed and embed in para.part.rels:
                            rel = para.part.rels[embed]
                            img_data = rel.target_part.blob
                            img_name = rel.target_part.partname.split("/")[-1]
                            image_count = _append_image(parts, image_count, img_data, img_name)
    return ExtractedContent(text="\n".join(text_parts), parts=parts)


def _extract_xlsx(content: bytes) -> ExtractedContent:
    from openpyxl import load_workbook
    wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    wb.close()
    text = "\n".join(lines)
    return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])


def _extract_pptx(content: bytes) -> ExtractedContent:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(BytesIO(content))
    text_parts: list[str] = []
    parts: list[dict] = []
    image_count = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_header = f"## Slide {slide_num}"
        text_parts.append(slide_header)
        parts.append({"type": "text", "text": slide_header})
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
                        parts.append({"type": "text", "text": text})
            if image_count < MAX_IMAGES and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_data = shape.image.blob
                img_name = shape.image.filename or "image.png"
                image_count = _append_image(parts, image_count, img_data, img_name)
    return ExtractedContent(text="\n".join(text_parts), parts=parts)


def _extract_html(content: bytes) -> ExtractedContent:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "lxml")
    text = soup.get_text(separator="\n")
    return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])


def _extract_plain_text(content: bytes) -> ExtractedContent:
    text = content.decode("utf-8", errors="replace")
    return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])


def extract_content(content: bytes, content_type: str, filename: str) -> ExtractedContent:
    """Extract text and content_parts from an uploaded document.

    Dispatches by content_type (falling back to filename extension) across
    PDF, DOCX, XLSX, PPTX, HTML, and plain text/markdown. Anything
    unrecognized is decoded as UTF-8 text.
    """
    if content_type == "application/pdf" or filename.endswith(".pdf"):
        return _extract_pdf(content)

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ) or filename.endswith(".docx"):
        return _extract_docx(content)

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ) or filename.endswith(".xlsx"):
        return _extract_xlsx(content)

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ) or filename.endswith(".pptx"):
        return _extract_pptx(content)

    elif content_type == "text/html" or filename.endswith(".html") or filename.endswith(".htm"):
        return _extract_html(content)

    elif content_type.startswith("text/") or filename.endswith((".txt", ".md")):
        return _extract_plain_text(content)

    return _extract_plain_text(content)
