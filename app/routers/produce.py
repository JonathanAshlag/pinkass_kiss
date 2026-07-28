"""Document ingestion endpoint."""

import base64
import logging
from dataclasses import dataclass, field
from io import BytesIO

from fastapi import APIRouter, Depends, Form, UploadFile, File, HTTPException

from app.container import PageRepo, RequestRepo, SourceFileRepo
from app.infrastructure.mongo import get_gridfs
from app.models.page import TrustTier
from app.models.user import User
from app.routers.deps import require_editor
from app.llm.pipeline import run_ingestion_pipeline

logger = logging.getLogger("pinkas.produce")

router = APIRouter(tags=["produce"])

MAX_IMAGES = 20


@dataclass
class ExtractedContent:
    text: str
    parts: list[dict] = field(default_factory=list)


def _encode_image(data: bytes, name: str) -> str:
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else "png"
    mime = f"image/{ext}" if ext in ("png", "jpeg", "jpg", "gif", "webp") else "image/png"
    b64 = base64.b64encode(data).decode("ascii")
    return f"data:{mime};base64,{b64}"


def _extract_content(content: bytes, content_type: str, filename: str) -> ExtractedContent:
    if content_type == "application/pdf" or filename.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(BytesIO(content))
        text_parts: list[str] = []
        parts: list[dict] = []
        image_count = 0
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
                parts.append({"type": "text", "text": page_text})
            if image_count < MAX_IMAGES:
                for img in page.images:
                    if image_count >= MAX_IMAGES:
                        break
                    uri = _encode_image(img.data, img.name)
                    parts.append({"type": "image_url", "image_url": {"url": uri}})
                    image_count += 1
        return ExtractedContent(text="\n".join(text_parts), parts=parts)

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ) or filename.endswith(".docx"):
        from docx import Document
        doc = Document(BytesIO(content))
        text_parts: list[str] = []
        parts: list[dict] = []
        image_count = 0
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                parts.append({"type": "text", "text": para.text})
            if image_count < MAX_IMAGES:
                for run in para.runs:
                    for drawing in run.element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                        blips = drawing.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                        for blip in blips:
                            if image_count >= MAX_IMAGES:
                                break
                            embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if embed and embed in para.part.rels:
                                rel = para.part.rels[embed]
                                img_data = rel.target_part.blob
                                img_name = rel.target_part.partname.split("/")[-1]
                                uri = _encode_image(img_data, img_name)
                                parts.append({"type": "image_url", "image_url": {"url": uri}})
                                image_count += 1
        return ExtractedContent(text="\n".join(text_parts), parts=parts)

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ) or filename.endswith(".xlsx"):
        from openpyxl import load_workbook
        wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append("\t".join(cells))
        wb.close()
        text = "\n".join(lines)
        return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])

    elif content_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ) or filename.endswith(".pptx"):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(BytesIO(content))
        text_parts: list[str] = []
        parts: list[dict] = []
        image_count = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_header = f"## Slide {slide_num}"
            text_parts.append(slide_header)
            parts.append({"type": "text", "text": slide_header})
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
                            parts.append({"type": "text", "text": text})
                if image_count < MAX_IMAGES and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_data = shape.image.blob
                    img_name = shape.image.filename or "image.png"
                    uri = _encode_image(img_data, img_name)
                    parts.append({"type": "image_url", "image_url": {"url": uri}})
                    image_count += 1
        return ExtractedContent(text="\n".join(text_parts), parts=parts)

    elif content_type == "text/html" or filename.endswith(".html") or filename.endswith(".htm"):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content, "lxml")
        text = soup.get_text(separator="\n")
        return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])

    elif content_type.startswith("text/") or filename.endswith((".txt", ".md")):
        text = content.decode("utf-8", errors="replace")
        return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])

    text = content.decode("utf-8", errors="replace")
    return ExtractedContent(text=text, parts=[{"type": "text", "text": text}])


@router.post("/produce")
async def produce(
    files: list[UploadFile] = File(...),
    initial_trust_tier: str = Form(default=TrustTier.unverified.value),
    user: User = Depends(require_editor),
    page_repo: PageRepo = None,
    req_repo: RequestRepo = None,
    source_file_repo: SourceFileRepo = None,
):
    """Ingest files, extract text, generate pages via LLM."""
    results = []

    for upload_file in files:
        content = await upload_file.read()
        content_type = upload_file.content_type or "application/octet-stream"
        filename = upload_file.filename or "unnamed"

        gridfs = get_gridfs()
        oid = await gridfs.upload_from_stream(
            filename,
            BytesIO(content),
            metadata={"content_type": content_type, "uploaded_by": user.user_id},
        )
        file_id_str = str(oid)

        extracted = _extract_content(content, content_type, filename)

        await source_file_repo.create(
            file_id=file_id_str,
            filename=filename,
            content_type=content_type,
            uploaded_by=user.user_id,
            extracted_text=extracted.text,
        )

        pipeline_results = await run_ingestion_pipeline(
            text=extracted.text,
            content_parts=extracted.parts,
            filename=filename,
            file_id_str=file_id_str,
            user=user,
            initial_trust_tier=initial_trust_tier,
            page_repo=page_repo,
            req_repo=req_repo,
        )

        generated_page_ids = [r["page_id"] for r in pipeline_results]
        results.extend(pipeline_results)

        await source_file_repo.set_page_ids(file_id_str, generated_page_ids)

    return {"generated_pages": results}


@router.get("/files/{file_id}")
async def get_file(file_id: str):
    """Download a stored source file."""
    from fastapi.responses import StreamingResponse

    from bson import ObjectId
    gridfs = get_gridfs()
    try:
        stream = await gridfs.open_download_stream(ObjectId(file_id))
    except Exception:
        raise HTTPException(status_code=404, detail="File not found")

    async def iterfile():
        while True:
            chunk = await stream.read(8192)
            if not chunk:
                break
            yield chunk

    return StreamingResponse(
        iterfile(),
        media_type=stream.metadata.get("content_type", "application/octet-stream") if stream.metadata else "application/octet-stream",
            headers={"Content-Disposition": f"attachment; filename={stream.filename}"},
        )
